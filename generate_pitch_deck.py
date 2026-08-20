"""
ClearSpend — 15-Slide Detailed Presentation Generator
Generates a professional PPTX pitch deck using python-pptx.
"""

import collections, collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand palette ──────────────────────────────────────────────────────
TEAL        = RGBColor(15, 118, 110)   # Primary accent
TEAL_LIGHT  = RGBColor(204, 251, 241)  # Light teal for backgrounds
TEAL_MID    = RGBColor(20, 184, 166)   # Mid teal
OFF_WHITE   = RGBColor(250, 250, 248)
DARK        = RGBColor(24, 24, 27)
GRAY        = RGBColor(113, 113, 122)
WHITE       = RGBColor(255, 255, 255)
ORANGE      = RGBColor(249, 115, 22)
GREEN       = RGBColor(16, 185, 129)
BLUE        = RGBColor(59, 130, 246)
RED         = RGBColor(239, 68, 68)
PURPLE      = RGBColor(139, 92, 246)

SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H


# ── Helpers ─────────────────────────────────────────────────────────────
def _set_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def _add_shape(slide, left, top, w, h, color, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def _add_textbox(slide, left, top, w, h, text, size=18, bold=False, color=DARK,
                 align=PP_ALIGN.LEFT, font_name="Calibri"):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf


def _add_bullets(tf, items, size=16, color=DARK, indent_color=GRAY, spacing=Pt(10)):
    """Add bullet items to an existing text-frame. Items starting with '→' are sub-bullets."""
    for item in items:
        p = tf.add_paragraph()
        is_sub = item.startswith("→")
        text = item.lstrip("→ ")
        p.text = text
        p.font.size = Pt(size - 2 if is_sub else size)
        p.font.color.rgb = indent_color if is_sub else color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 1 if is_sub else 0


def _accent_bar(slide, top=Inches(0)):
    """Full-width teal accent bar at top."""
    _add_shape(slide, Inches(0), top, SLIDE_W, Inches(0.06), TEAL)


def _slide_number(slide, num):
    _add_textbox(slide, Inches(9.2), Inches(7.05), Inches(0.6), Inches(0.35),
                 str(num), size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def blank(bg=OFF_WHITE):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    _set_bg(s, bg)
    return s


# ════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════
s = blank(TEAL)
_add_shape(s, Inches(0), Inches(0), SLIDE_W, Inches(0.12), TEAL_MID)

_add_textbox(s, Inches(0.8), Inches(1.8), Inches(8.4), Inches(1.5),
             "ClearSpend", size=60, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
_add_textbox(s, Inches(0.8), Inches(3.2), Inches(8.4), Inches(1.0),
             "Autonomous AI-Powered Expense Tracking\n& Financial Coach for India",
             size=26, color=TEAL_LIGHT, align=PP_ALIGN.LEFT)

_add_shape(s, Inches(0.8), Inches(4.6), Inches(1.2), Inches(0.05), TEAL_LIGHT)

_add_textbox(s, Inches(0.8), Inches(5.0), Inches(8.4), Inches(0.8),
             "Built with React 18 · TypeScript · Supabase · Gemini AI\n"
             "clearspend-ai-expense-tracker.vercel.app",
             size=14, color=TEAL_LIGHT, align=PP_ALIGN.LEFT)
_slide_number(s, 1)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ════════════════════════════════════════════════════════════════════════
s = blank()
_accent_bar(s)
_add_textbox(s, Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.8),
             "The Problem", size=36, bold=True, color=TEAL)
_add_textbox(s, Inches(0.8), Inches(1.2), Inches(8.4), Inches(0.5),
             "Why 73% of Indians abandon expense trackers within 2 weeks",
             size=16, color=GRAY, bold=True)

problems = [
    ("😩", "Too Many Taps", "Traditional apps need 6+ taps per transaction — select amount, category, wallet, date, merchant, and confirm."),
    ("🤖", "Zero Intelligence", "They just record data. No pattern recognition, no proactive alerts, no learning from your behavior."),
    ("📊", "Reactive, Not Predictive", "You only discover you overspent on food after the month ends. By then, the money is gone."),
    ("🔁", "Duplicate Chaos", "UPI failures cause double-debits. No app catches them — your ledger silently drifts from reality."),
]

for i, (emoji, title, desc) in enumerate(problems):
    y = Inches(1.9 + i * 1.25)
    _add_shape(s, Inches(0.8), y, Inches(0.7), Inches(0.7), TEAL_LIGHT, MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_textbox(s, Inches(0.85), y + Inches(0.08), Inches(0.6), Inches(0.55),
                 emoji, size=24, align=PP_ALIGN.CENTER)
    _add_textbox(s, Inches(1.7), y + Inches(0.02), Inches(7.5), Inches(0.35),
                 title, size=18, bold=True, color=DARK)
    _add_textbox(s, Inches(1.7), y + Inches(0.38), Inches(7.5), Inches(0.5),
                 desc, size=13, color=GRAY)
_slide_number(s, 2)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Solution
# ════════════════════════════════════════════════════════════════════════
s = blank()
_accent_bar(s)
_add_textbox(s, Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.8),
             "The Solution: ClearSpend", size=36, bold=True, color=TEAL)
_add_textbox(s, Inches(0.8), Inches(1.2), Inches(8.4), Inches(0.5),
             "An AI-first financial copilot that thinks like you",
             size=16, color=GRAY, bold=True)

solutions = [
    ("Natural Language Input", "Type how you think: '380 zomato lunch' — done in 3 seconds.", ORANGE),
    ("Autonomous AI Parsing", "Extracts amount (₹380), merchant (Zomato), category (Food), date — with confidence scores.", GREEN),
    ("Self-Learning Rules", "Corrects learn. Change 'Groceries' → 'Supplies' once, and ClearSpend remembers forever.", BLUE),
    ("Proactive Budget Coach", "Daily pace alerts: 'At this pace you'll spend ₹15,000 on Food. Cap ₹300/day.'", PURPLE),
]

for i, (title, desc, accent) in enumerate(solutions):
    y = Inches(1.9 + i * 1.3)
    _add_shape(s, Inches(0.8), y, Inches(0.08), Inches(0.9), accent)
    _add_textbox(s, Inches(1.1), y + Inches(0.02), Inches(8.1), Inches(0.35),
                 title, size=18, bold=True, color=DARK)
    _add_textbox(s, Inches(1.1), y + Inches(0.4), Inches(8.1), Inches(0.5),
                 desc, size=14, color=GRAY)
_slide_number(s, 3)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 4 — How It Works: Quick Add Flow
# ════════════════════════════════════════════════════════════════════════
s = blank()
_accent_bar(s)
_add_textbox(s, Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.8),
             "How It Works: The Quick Add Flow", size=36, bold=True, color=TEAL)

steps = [
    ("1", "User Types", '"2k rent yesterday"', TEAL),
    ("2", "NLP Parses", "Amount: ₹2,000\nCategory: Rent\nDate: Yesterday", ORANGE),
    ("3", "Confidence Card", "Shows parsed result\nwith High / Check / Guess\nbadges", GREEN),
    ("4", "One-Tap Save", "Transaction saved\noptimistically with\nskeleton loading", BLUE),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(0.6 + i * 2.35)
    y = Inches(1.8)
    # Circle number
    circle = _add_shape(s, x + Inches(0.65), y, Inches(0.6), Inches(0.6), color, MSO_SHAPE.OVAL)
    circle.text_frame.paragraphs[0].text = num
    circle.text_frame.paragraphs[0].font.size = Pt(22)
    circle.text_frame.paragraphs[0].font.bold = True
    circle.text_frame.paragraphs[0].font.color.rgb = WHITE
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Arrow (except last)
    if i < 3:
        _add_textbox(s, x + Inches(1.9), y + Inches(0.1), Inches(0.5), Inches(0.4),
                     "→", size=28, color=GRAY, align=PP_ALIGN.CENTER)
    # Title
    _add_textbox(s, x, y + Inches(0.8), Inches(2.0), Inches(0.4),
                 title, size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    # Description
    _add_textbox(s, x, y + Inches(1.2), Inches(2.0), Inches(1.2),
                 desc, size=12, color=GRAY, align=PP_ALIGN.CENTER)

# Indian shorthand examples box
_add_shape(s, Inches(0.8), Inches(4.8), Inches(8.4), Inches(2.0), TEAL_LIGHT, MSO_SHAPE.ROUNDED_RECTANGLE)
_add_textbox(s, Inches(1.0), Inches(4.9), Inches(8.0), Inches(0.4),
             "Indian Shorthand Examples Supported:", size=14, bold=True, color=TEAL)

examples = [
    '"380 zomato lunch"  →  ₹380 · Food & Dining · Today',
    '"2k rent yesterday"  →  ₹2,000 · Rent · Aug 20',
    '"got 50000 salary"  →  ₹50,000 · Salary · Income',
    '"1.5k uber cab"  →  ₹1,500 · Transport · Today',
]
tf = _add_textbox(s, Inches(1.0), Inches(5.3), Inches(8.0), Inches(1.4),
                  examples[0], size=12, color=DARK)
tf.paragraphs[0].font.name = "Consolas"
for ex in examples[1:]:
    p = tf.add_paragraph()
    p.text = ex
    p.font.size = Pt(12)
    p.font.color.rgb = DARK
    p.font.name = "Consolas"
    p.space_after = Pt(4)
_slide_number(s, 4)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 5 — AI Confidence & Trust
# ════════════════════════════════════════════════════════════════════════
s = blank()
_accent_bar(s)
_add_textbox(s, Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.8),
             "AI Confidence & Transparency", size=36, bold=True, color=TEAL)
_add_textbox(s, Inches(0.8), Inches(1.2), Inches(8.4), Inches(0.5),
             "Every AI decision is explainable — the user always stays in control",
             size=16, color=GRAY, bold=True)

badges = [
    ("HIGH CONFIDENCE", "≥ 80%", "AI is very sure. One-tap save.", GREEN, "Examples: 'zomato' → Food, 'uber' → Transport. Matched via category rules or strong LLM signal."),
    ("CHECK THIS", "50–79%", "AI has a reasonable guess but wants human confirmation.", ORANGE, "Examples: 'amazon 2k' could be Shopping OR Groceries. User picks and the system learns."),
    ("GUESS", "< 50%", "Low confidence. Falls back to 'Other' category.", RED, "Examples: Ambiguous input like '500 misc'. User corrects and a new rule is created."),
]

for i, (label, pct, desc, color, detail) in enumerate(badges):
    y = Inches(2.0 + i * 1.7)
    _add_shape(s, Inches(0.8), y, Inches(2.0), Inches(0.5), color, MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_textbox(s, Inches(0.85), y + Inches(0.05), Inches(1.9), Inches(0.4),
                 label, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(s, Inches(3.0), y + Inches(0.0), Inches(1.0), Inches(0.4),
                 pct, size=16, bold=True, color=color, align=PP_ALIGN.CENTER)
    _add_textbox(s, Inches(4.0), y + Inches(0.0), Inches(5.5), Inches(0.4),
                 desc, size=14, bold=True, color=DARK)
    _add_textbox(s, Inches(4.0), y + Inches(0.45), Inches(5.5), Inches(0.7),
                 detail, size=12, color=GRAY)
_slide_number(s, 5)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Self-Learning Rule Engine
# ════════════════════════════════════════════════════════════════════════
s = blank()
_accent_bar(s)
_add_textbox(s, Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.8),
             "Self-Learning Category Rules", size=36, bold=True, color=TEAL)
_add_textbox(s, Inches(0.8), Inches(1.2), Inches(8.4), Inches(0.5),
             "ClearSpend gets smarter with every correction you make",
             size=16, color=GRAY, bold=True)

flow = [
    "① AI parses '450 dmart monthly supplies' → Suggests 'Groceries'",
    "② User changes category to 'Shopping'",
    "③ Engine upserts rule: 'dmart' → Shopping (hit_count: 1)",
    "④ Toast appears: 'Learned! Also apply to past transactions?'",
    "⑤ User taps 'Yes' → Bulk retroactive update on all 'dmart' records",
    "⑥ Next time 'dmart' appears → Auto-categorized as Shopping (100% confidence)",
]

for i, step in enumerate(flow):
    y = Inches(1.9 + i * 0.8)
    _add_shape(s, Inches(0.8), y, Inches(8.4), Inches(0.65), TEAL_LIGHT if i % 2 == 0 else OFF_WHITE,
               MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_textbox(s, Inches(1.0), y + Inches(0.1), Inches(8.0), Inches(0.45),
                 step, size=15, color=DARK)

_add_textbox(s, Inches(0.8), Inches(6.8), Inches(8.4), Inches(0.4),
             "Result: Zero-shot learning → The more you use ClearSpend, the less you need to correct it.",
             size=14, bold=True, color=TEAL)
_slide_number(s, 6)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Review Inbox: Duplicate & Anomaly Detection
# ════════════════════════════════════════════════════════════════════════
s = blank()
_accent_bar(s)
_add_textbox(s, Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.8),
             "Review Inbox: Your Financial Guardian", size=36, bold=True, color=TEAL)

# Left column — Duplicates
_add_shape(s, Inches(0.6), Inches(1.5), Inches(4.2), Inches(5.2), TEAL_LIGHT, MSO_SHAPE.ROUNDED_RECTANGLE)
_add_textbox(s, Inches(0.8), Inches(1.7), Inches(3.8), Inches(0.4),
             "🔁 Duplicate Detection", size=20, bold=True, color=TEAL)

dup_items = [
    "Algorithm: Token Jaccard + Levenshtein distance",
    "Fingerprint: hash(merchant + amount + date)",
    "Catches UPI double-debits & repeated charges",
    "Side-by-side comparison view",
    "Actions: Merge (Keep Older) · Keep Both · Delete",
]
tf = _add_textbox(s, Inches(0.8), Inches(2.3), Inches(3.8), Inches(4.0),
                  dup_items[0], size=13, color=DARK)
for item in dup_items[1:]:
    p = tf.add_paragraph()
    p.text = "• " + item
    p.font.size = Pt(13)
    p.font.color.rgb = DARK
    p.space_after = Pt(10)
tf.paragraphs[0].text = "• " + dup_items[0]

# Right column — Anomalies
_add_shape(s, Inches(5.2), Inches(1.5), Inches(4.2), Inches(5.2), TEAL_LIGHT, MSO_SHAPE.ROUNDED_RECTANGLE)
_add_textbox(s, Inches(5.4), Inches(1.7), Inches(3.8), Inches(0.4),
             "⚠️ Anomaly Detection", size=20, bold=True, color=ORANGE)

anom_items = [
    "Threshold: > 3× the category's median",
    "Example: ₹4,500 on Food when median is ₹380",
    "Highlights unusual spending patterns",
    "Quick 'Confirm Correct' or 'Recategorize'",
    "Prevents silent budget leaks",
]
tf2 = _add_textbox(s, Inches(5.4), Inches(2.3), Inches(3.8), Inches(4.0),
                   anom_items[0], size=13, color=DARK)
for item in anom_items[1:]:
    p = tf2.add_paragraph()
    p.text = "• " + item
    p.font.size = Pt(13)
    p.font.color.rgb = DARK
    p.space_after = Pt(10)
tf2.paragraphs[0].text = "• " + anom_items[0]
_slide_number(s, 7)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Proactive Budget Forecasting
# ════════════════════════════════════════════════════════════════════════
s = blank()
_accent_bar(s)
_add_textbox(s, Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.8),
             "Proactive Budget Forecasting", size=36, bold=True, color=TEAL)
_add_textbox(s, Inches(0.8), Inches(1.2), Inches(8.4), Inches(0.5),
             "From reactive tracking → predictive financial coaching",
             size=16, color=GRAY, bold=True)

# Formula box
_add_shape(s, Inches(0.8), Inches(1.9), Inches(8.4), Inches(1.2), TEAL, MSO_SHAPE.ROUNDED_RECTANGLE)
_add_textbox(s, Inches(1.0), Inches(2.0), Inches(8.0), Inches(0.5),
             "THE PACING FORMULA", size=14, bold=True, color=TEAL_LIGHT)
_add_textbox(s, Inches(1.0), Inches(2.4), Inches(8.0), Inches(0.5),
             "Projected Spend = (Spent ÷ Elapsed Days) × Days in Month",
             size=22, bold=True, color=WHITE, font_name="Consolas")

# Example scenario
_add_textbox(s, Inches(0.8), Inches(3.5), Inches(8.4), Inches(0.4),
             "Real Example (Day 15 of a 31-day month):", size=16, bold=True, color=DARK)

scenario = [
    "Food Budget: ₹9,000/month   |   Spent so far: ₹5,800",
    "Daily Pace: ₹5,800 ÷ 15 = ₹387/day",
    "Projected: ₹387 × 31 = ₹11,987  (133% of budget!)",
    "",
    "⚡ Alert Banner:",
    '"At this pace you\'ll spend ₹11,987 on Food & Dining."',
    '"Cap it at ₹200/day for the remaining 16 days to stay on track."',
]
tf = _add_textbox(s, Inches(1.0), Inches(3.9), Inches(8.0), Inches(3.0),
                  scenario[0], size=14, color=DARK)
for line in scenario[1:]:
    p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(14)
    p.font.color.rgb = TEAL if line.startswith('"') else (ORANGE if line.startswith("⚡") else DARK)
    p.font.bold = line.startswith("⚡") or line.startswith('"')
    p.font.name = "Calibri"
    p.space_after = Pt(6)
_slide_number(s, 8)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Dashboard & Data Visualization
# ════════════════════════════════════════════════════════════════════════
s = blank()
_accent_bar(s)
_add_textbox(s, Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.8),
             "Dashboard & Analytics", size=36, bold=True, color=TEAL)

features = [
    ("📊", "Interactive Donut Chart", "Recharts-powered category breakdown with hover tooltips. Tap any slice to filter the transaction list."),
    ("💰", "Summary Cards", "Real-time totals for Income, Expenses, and Balance. Updated optimistically on every quick-add."),
    ("📋", "Day-Grouped Ledger", "Transactions sorted by date with sticky day headers. Multi-select mode for bulk recategorize or delete."),
    ("🔍", "Smart Filtering", "Filter by category, wallet, date range, or keyword. Combine filters for laser-focused analysis."),
    ("📁", "CSV Export", "One-click export of all transactions to CSV for spreadsheets, accountants, or tax filing."),
    ("🏆", "AI Insights Tab", "Finance coach cards: top movers, subscription detection, spending streaks, forecast summaries."),
]

for i, (emoji, title, desc) in enumerate(features):
    col = 0 if i < 3 else 1
    row = i % 3
    x = Inches(0.6 + col * 4.6)
    y = Inches(1.6 + row * 1.8)
    _add_shape(s, x, y, Inches(4.3), Inches(1.5), TEAL_LIGHT, MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_textbox(s, x + Inches(0.15), y + Inches(0.1), Inches(4.0), Inches(0.4),
                 f"{emoji}  {title}", size=16, bold=True, color=TEAL)
    _add_textbox(s, x + Inches(0.15), y + Inches(0.55), Inches(4.0), Inches(0.85),
                 desc, size=12, color=GRAY)
_slide_number(s, 9)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Tech Stack Architecture
# ════════════════════════════════════════════════════════════════════════
s = blank()
_accent_bar(s)
_add_textbox(s, Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.8),
             "Technology Architecture", size=36, bold=True, color=TEAL)

# Frontend layer
_add_shape(s, Inches(0.6), Inches(1.5), Inches(4.0), Inches(2.6), TEAL_LIGHT, MSO_SHAPE.ROUNDED_RECTANGLE)
_add_textbox(s, Inches(0.8), Inches(1.6), Inches(3.6), Inches(0.4),
             "⚛️  FRONTEND", size=18, bold=True, color=TEAL)
tf = _add_textbox(s, Inches(0.8), Inches(2.1), Inches(3.6), Inches(1.8),
                  "• React 18 + TypeScript", size=13, color=DARK)
for t in ["• Vite (lightning-fast HMR)", "• Tailwind CSS (design tokens)",
          "• Recharts (data visualization)", "• Lucide React (icon system)",
          "• Context + localStorage state"]:
    p = tf.add_paragraph(); p.text = t; p.font.size = Pt(13); p.font.color.rgb = DARK; p.space_after = Pt(4)

# Backend layer
_add_shape(s, Inches(5.4), Inches(1.5), Inches(4.0), Inches(2.6), TEAL_LIGHT, MSO_SHAPE.ROUNDED_RECTANGLE)
_add_textbox(s, Inches(5.6), Inches(1.6), Inches(3.6), Inches(0.4),
             "🗄️  BACKEND", size=18, bold=True, color=TEAL)
tf = _add_textbox(s, Inches(5.6), Inches(2.1), Inches(3.6), Inches(1.8),
                  "• Supabase (PostgreSQL)", size=13, color=DARK)
for t in ["• Row-Level Security (RLS)", "• Deno Edge Functions",
          "• Auth (Magic Link / OAuth)", "• 7 relational tables", "• Auto-seed triggers"]:
    p = tf.add_paragraph(); p.text = t; p.font.size = Pt(13); p.font.color.rgb = DARK; p.space_after = Pt(4)

# AI layer
_add_shape(s, Inches(0.6), Inches(4.4), Inches(8.8), Inches(2.4), TEAL, MSO_SHAPE.ROUNDED_RECTANGLE)
_add_textbox(s, Inches(0.8), Inches(4.5), Inches(8.4), Inches(0.4),
             "🤖  AI / AGENTIC LAYER", size=18, bold=True, color=WHITE)
tf = _add_textbox(s, Inches(0.8), Inches(5.0), Inches(4.0), Inches(1.5),
                  "• Model-agnostic LLM interface", size=13, color=TEAL_LIGHT)
for t in ["• Gemini 2.5 Flash (primary)", "• GPT-4o-mini (fallback)",
          "• Deterministic regex fallback"]:
    p = tf.add_paragraph(); p.text = t; p.font.size = Pt(13); p.font.color.rgb = TEAL_LIGHT; p.space_after = Pt(4)

tf2 = _add_textbox(s, Inches(5.4), Inches(5.0), Inches(4.0), Inches(1.5),
                   "Edge Functions:", size=13, bold=True, color=WHITE)
for t in ["• /parse-transaction", "• /detect-duplicates",
          "• /generate-insights"]:
    p = tf2.add_paragraph(); p.text = t; p.font.size = Pt(13); p.font.color.rgb = TEAL_LIGHT; p.font.name = "Consolas"; p.space_after = Pt(4)
_slide_number(s, 10)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Database Schema
# ════════════════════════════════════════════════════════════════════════
s = blank()
_accent_bar(s)
_add_textbox(s, Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.8),
             "Database Schema (7 Tables)", size=36, bold=True, color=TEAL)

tables = [
    ("profiles", "User identity, display name, base currency (INR), onboarding status"),
    ("wallets", "HDFC Bank, Google Pay, Cash, ICICI Card — with opening balances"),
    ("categories", "12 default categories (10 expense + 2 income), custom icons & colors"),
    ("transactions", "Core ledger: amount, kind, merchant, AI confidence, fingerprint, status"),
    ("budgets", "Monthly category limits with alert thresholds (default 80%)"),
    ("category_rules", "Self-learned merchant→category mappings with hit counts"),
    ("insights", "AI-generated finance coach cards: forecasts, top movers, streaks"),
]

for i, (name, desc) in enumerate(tables):
    y = Inches(1.5 + i * 0.78)
    _add_shape(s, Inches(0.8), y, Inches(2.0), Inches(0.55), TEAL, MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_textbox(s, Inches(0.85), y + Inches(0.07), Inches(1.9), Inches(0.4),
                 name, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Consolas")
    _add_textbox(s, Inches(3.0), y + Inches(0.07), Inches(6.5), Inches(0.45),
                 desc, size=13, color=DARK)

_add_textbox(s, Inches(0.8), Inches(7.0), Inches(8.4), Inches(0.35),
             "All tables protected by Row Level Security: user_id = auth.uid()",
             size=12, bold=True, color=TEAL)
_slide_number(s, 11)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Design Philosophy & Aesthetics
# ════════════════════════════════════════════════════════════════════════
s = blank()
_accent_bar(s)
_add_textbox(s, Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.8),
             "Design Philosophy", size=36, bold=True, color=TEAL)
_add_textbox(s, Inches(0.8), Inches(1.2), Inches(8.4), Inches(0.5),
             "Financial apps should feel calm, trustworthy, and premium",
             size=16, color=GRAY, bold=True)

# Color palette swatches
colors_data = [
    ("#0F766E", "Deep Teal", "Primary accent — trust, calm", TEAL),
    ("#FAFAF8", "Off-White", "Canvas background — breathable", RGBColor(230, 230, 228)),
    ("#18181B", "Near-Black", "Body text — clarity", DARK),
    ("#F97316", "Warm Orange", "Food, alerts, warnings", ORANGE),
    ("#10B981", "Emerald", "Income, success states", GREEN),
    ("#8B5CF6", "Violet", "Entertainment, insights", PURPLE),
]

for i, (hex_val, name, purpose, color) in enumerate(colors_data):
    x = Inches(0.6 + (i % 3) * 3.1)
    y = Inches(1.9 + (i // 3) * 1.6)
    _add_shape(s, x, y, Inches(0.6), Inches(0.6), color, MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_textbox(s, x + Inches(0.75), y + Inches(0.0), Inches(2.1), Inches(0.3),
                 f"{name}  {hex_val}", size=13, bold=True, color=DARK)
    _add_textbox(s, x + Inches(0.75), y + Inches(0.3), Inches(2.1), Inches(0.3),
                 purpose, size=11, color=GRAY)

# Design principles
_add_textbox(s, Inches(0.8), Inches(5.2), Inches(8.4), Inches(0.4),
             "Key Design Principles:", size=16, bold=True, color=TEAL)
principles = [
    "Tabular numerals (font-variant-numeric: tabular-nums) for perfect number alignment",
    "Skeleton loaders during AI parsing — the UI feels alive, never frozen",
    "Optimistic updates — transactions appear instantly before server confirms",
    "Micro-animations on confidence badges, toast notifications, and card transitions",
]
tf = _add_textbox(s, Inches(0.8), Inches(5.7), Inches(8.4), Inches(1.5),
                  "• " + principles[0], size=13, color=DARK)
for pr in principles[1:]:
    p = tf.add_paragraph()
    p.text = "• " + pr
    p.font.size = Pt(13)
    p.font.color.rgb = DARK
    p.space_after = Pt(6)
_slide_number(s, 12)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Demo Mode & Onboarding
# ════════════════════════════════════════════════════════════════════════
s = blank()
_accent_bar(s)
_add_textbox(s, Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.8),
             "Instant Demo & Onboarding", size=36, bold=True, color=TEAL)

# Left: Demo Mode
_add_shape(s, Inches(0.6), Inches(1.5), Inches(4.2), Inches(5.2), TEAL_LIGHT, MSO_SHAPE.ROUNDED_RECTANGLE)
_add_textbox(s, Inches(0.8), Inches(1.6), Inches(3.8), Inches(0.4),
             "🎮  1-Click Demo Mode", size=18, bold=True, color=TEAL)
demo_items = [
    "40+ realistic pre-seeded transactions",
    "2 months of spending data (July & Aug 2026)",
    "4 wallets: HDFC, Google Pay, Cash, ICICI",
    "12 categories with custom colors & icons",
    "Pre-planted duplicate pair for review",
    "Pre-planted anomaly (₹4,500 > 3× median)",
    "3 budget targets with pace calculations",
    "No Supabase credentials needed!",
]
tf = _add_textbox(s, Inches(0.8), Inches(2.2), Inches(3.8), Inches(4.0),
                  "• " + demo_items[0], size=13, color=DARK)
for item in demo_items[1:]:
    p = tf.add_paragraph()
    p.text = "• " + item
    p.font.size = Pt(13)
    p.font.color.rgb = DARK
    p.space_after = Pt(8)

# Right: Onboarding
_add_shape(s, Inches(5.2), Inches(1.5), Inches(4.2), Inches(5.2), TEAL_LIGHT, MSO_SHAPE.ROUNDED_RECTANGLE)
_add_textbox(s, Inches(5.4), Inches(1.6), Inches(3.8), Inches(0.4),
             "🚀  3-Step Onboarding Wizard", size=18, bold=True, color=TEAL)
onboard = [
    "Step 1: Welcome & Profile Setup",
    "→ Display name, preferred currency",
    "Step 2: Create Your First Wallet",
    "→ Bank, UPI, Cash, or Credit Card",
    "Step 3: Set Monthly Budgets",
    "→ 3-month average suggestions shown",
    "",
    "Goal: User reaches 'Aha!' moment",
    "within 60 seconds of landing.",
]
tf = _add_textbox(s, Inches(5.4), Inches(2.2), Inches(3.8), Inches(4.0),
                  onboard[0], size=13, bold=True, color=DARK)
for item in onboard[1:]:
    p = tf.add_paragraph()
    is_sub = item.startswith("→")
    p.text = item
    p.font.size = Pt(12 if is_sub else 13)
    p.font.color.rgb = GRAY if is_sub else DARK
    p.font.bold = not is_sub and item != ""
    p.space_after = Pt(6)
_slide_number(s, 13)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Roadmap & Future Vision
# ════════════════════════════════════════════════════════════════════════
s = blank()
_accent_bar(s)
_add_textbox(s, Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.8),
             "Roadmap & Future Vision", size=36, bold=True, color=TEAL)

phases = [
    ("NOW  ✅", "Hackathon MVP", [
        "NLP Quick Add with Indian shorthand",
        "Self-learning category rules",
        "Duplicate & anomaly detection",
        "Proactive budget pacing",
        "Deployed on Vercel + Supabase"
    ], GREEN),
    ("NEXT  🔜", "Phase 2", [
        "SMS webhook auto-capture",
        "Bank statement PDF import",
        "Multi-currency with live FX rates",
        "Push notification alerts",
        "Recurring transaction detection"
    ], ORANGE),
    ("VISION  🚀", "Phase 3", [
        "Investment portfolio tracking",
        "Tax-saving recommendations",
        "Family/shared expense groups",
        "WhatsApp bot integration",
        "AI financial planning assistant"
    ], PURPLE),
]

for i, (label, title, items, color) in enumerate(phases):
    x = Inches(0.5 + i * 3.15)
    _add_shape(s, x, Inches(1.5), Inches(2.95), Inches(0.55), color, MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_textbox(s, x + Inches(0.1), Inches(1.55), Inches(2.75), Inches(0.45),
                 f"{label}", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(s, x + Inches(0.1), Inches(2.15), Inches(2.75), Inches(0.35),
                 title, size=15, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    tf = _add_textbox(s, x + Inches(0.1), Inches(2.6), Inches(2.75), Inches(4.0),
                      "• " + items[0], size=12, color=DARK)
    for item in items[1:]:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(12)
        p.font.color.rgb = DARK
        p.space_after = Pt(6)
_slide_number(s, 14)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Thank You / CTA
# ════════════════════════════════════════════════════════════════════════
s = blank(TEAL)
_add_shape(s, Inches(0), Inches(0), SLIDE_W, Inches(0.12), TEAL_MID)

_add_textbox(s, Inches(0.8), Inches(1.5), Inches(8.4), Inches(1.0),
             "Thank You", size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
_add_textbox(s, Inches(0.8), Inches(2.8), Inches(8.4), Inches(0.6),
             "ClearSpend — Your AI Financial Copilot",
             size=24, color=TEAL_LIGHT, align=PP_ALIGN.CENTER)

_add_shape(s, Inches(3.5), Inches(3.8), Inches(3.0), Inches(0.04), TEAL_LIGHT)

_add_textbox(s, Inches(1.0), Inches(4.2), Inches(8.0), Inches(0.5),
             "🌐  Live App", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
_add_textbox(s, Inches(1.0), Inches(4.7), Inches(8.0), Inches(0.4),
             "clearspend-ai-expense-tracker.vercel.app",
             size=18, color=TEAL_LIGHT, align=PP_ALIGN.CENTER, font_name="Consolas")

_add_textbox(s, Inches(1.0), Inches(5.4), Inches(8.0), Inches(0.5),
             "📦  Source Code", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
_add_textbox(s, Inches(1.0), Inches(5.9), Inches(8.0), Inches(0.4),
             "github.com/sreetejlakkam/ClearSpend-AI-Expense-Tracker",
             size=16, color=TEAL_LIGHT, align=PP_ALIGN.CENTER, font_name="Consolas")

_add_textbox(s, Inches(1.0), Inches(6.7), Inches(8.0), Inches(0.5),
             "Built with ❤️ using React · TypeScript · Supabase · Gemini AI",
             size=14, color=TEAL_LIGHT, align=PP_ALIGN.CENTER)
_slide_number(s, 15)


# ── Save ───────────────────────────────────────────────────────────────
output_path = "/media/sreetej/SSD1/Antigravity/Budget_tracker/ClearSpend_Pitch_Deck.pptx"
prs.save(output_path)
print(f"✅ Presentation saved: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
